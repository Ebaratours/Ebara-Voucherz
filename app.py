
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
import uuid
import shutil
import os
from PIL import Image
from fpdf import FPDF


def generate_serial():
    return "VCH" + datetime.now().strftime('%y%m%d') + str(uuid.uuid4().hex)[:4]


def fill_voucher_template(name, date, serial_number):
    template_path = "Ebara Reservation-2025.pptx"
    output_path = f"voucher_{uuid.uuid4().hex}.pptx"
    shutil.copy(template_path, output_path)

    prs = Presentation(output_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if "M Mohamed Sareebu" in text:
                    shape.text_frame.text = text.replace("M Mohamed Sareebu", name)
                if "Date :" in text:
                    shape.text_frame.text = text.replace("Date :", f"Date : {date}")
                if "Res # :                    2501" in text:
                    shape.text_frame.text = text.replace("Res # :                    2501", f"Res # :                    {serial_number}")
    prs.save(output_path)
    return output_path


def convert_slide_to_pdf(pptx_path):
    img_dir = f"img_{uuid.uuid4().hex}"
    os.makedirs(img_dir, exist_ok=True)

    img_path = os.path.join(img_dir, "slide1.png")
    img = Image.new("RGB", (1280, 720), color=(255, 255, 255))
    img.save(img_path)

    pdf_path = pptx_path.replace(".pptx", ".pdf")
    pdf = FPDF()
    pdf.add_page()
    pdf.image(img_path, x=10, y=20, w=pdf.w - 20)
    pdf.output(pdf_path)

    shutil.rmtree(img_dir)
    return pdf_path


st.title("Ebara Voucher Generator")

with st.form("voucher_form"):
    guest_name = st.text_input("Guest Name")
    reservation_date = st.date_input("Reservation Date")
    submitted = st.form_submit_button("Generate Voucher")

if submitted:
    serial = generate_serial()
    pptx_path = fill_voucher_template(guest_name, reservation_date.strftime('%Y-%m-%d'), serial)
    pdf_path = convert_slide_to_pdf(pptx_path)

    with open(pdf_path, "rb") as f:
        st.success(f"Voucher generated for {guest_name} with Serial #{serial}")
        st.download_button(label="Download Voucher (PDF)", data=f, file_name="voucher.pdf", mime="application/pdf")

    os.remove(pptx_path)
    os.remove(pdf_path)
